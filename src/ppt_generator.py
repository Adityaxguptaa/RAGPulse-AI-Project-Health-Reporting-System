"""PowerPoint presentation generator — portfolio (7 slides) + per-project (5 slides)."""

from __future__ import annotations

import io
import logging
from datetime import date
from pathlib import Path
from typing import Any

from src.models import Project

logger = logging.getLogger(__name__)

# ── Palette ───────────────────────────────────────────────────────────────────
_GREEN = (46, 204, 113)
_AMBER = (243, 156, 18)
_RED = (231, 76, 60)
_BLUE = (52, 152, 219)
_DARK = (44, 62, 80)
_WHITE = (255, 255, 255)
_LIGHT = (248, 249, 250)
_GREY = (149, 165, 166)
_NAVY = (26, 37, 47)

_RAG_RGB = {"Green": _GREEN, "Amber": _AMBER, "Red": _RED, "Unknown": _GREY}


def _rgb(r, g, b):
    from pptx.dml.color import RGBColor

    return RGBColor(r, g, b)


def _set_bg(slide, r, g, b):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(r, g, b)


def _box(slide, left, top, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    s = slide.shapes.add_shape(1, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = _rgb(*fill)
    if line:
        s.line.color.rgb = _rgb(*line)
    else:
        s.line.fill.background()
    return s


def _text(
    slide,
    left,
    top,
    w,
    h,
    txt,
    size=12,
    bold=False,
    color=_WHITE,
    center=False,
    wrap=True,
):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = str(txt)
    if center:
        p.alignment = PP_ALIGN.CENTER
    if p.runs:
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)


def _rag_badge(slide, left, top, w, h, status):
    rgb = _RAG_RGB.get(status, _GREY)
    _box(slide, left, top, w, h, fill=rgb)
    _text(slide, left, top, w, h, status, size=11, bold=True, color=_WHITE, center=True)


def _new_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _prs():
    from pptx import Presentation
    from pptx.util import Inches

    p = Presentation()
    p.slide_width = Inches(13.33)
    p.slide_height = Inches(7.5)
    return p


# ── Shared header/footer helpers ──────────────────────────────────────────────


def _header_bar(slide, prs, title: str, bg=_DARK):
    from pptx.util import Inches

    W = prs.slide_width
    _box(slide, 0, 0, W, Inches(0.75), fill=bg)
    _text(
        slide,
        Inches(0.3),
        Inches(0.1),
        W - Inches(0.6),
        Inches(0.55),
        title,
        size=20,
        bold=True,
        color=_WHITE,
    )


def _footer_bar(slide, prs, left_txt: str, right_txt: str, bg=_DARK):
    from pptx.util import Inches

    W = prs.slide_width
    H = prs.slide_height
    _box(slide, 0, H - Inches(0.45), W, Inches(0.45), fill=bg)
    _text(
        slide,
        Inches(0.3),
        H - Inches(0.42),
        W * 0.5,
        Inches(0.38),
        left_txt,
        size=9,
        color=_WHITE,
    )
    _text(
        slide,
        W * 0.5,
        H - Inches(0.42),
        W * 0.5 - Inches(0.3),
        Inches(0.38),
        right_txt,
        size=9,
        color=_WHITE,
        center=True,
    )


# ══════════════════════════════════════════════════════════════════════════════
# PORTFOLIO SLIDES (7 slides)
# ══════════════════════════════════════════════════════════════════════════════


def _p_title(prs, summary):
    from pptx.util import Inches

    s = _new_blank(prs)
    W, H = prs.slide_width, prs.slide_height
    _set_bg(s, *_DARK)
    _box(s, 0, int(H * 0.55), W, int(H * 0.06), fill=_BLUE)
    _box(s, 0, H - Inches(0.55), W, Inches(0.55), fill=_NAVY)
    _text(
        s,
        Inches(1),
        Inches(1.2),
        W - Inches(2),
        Inches(0.6),
        "Monthly Portfolio Health Report",
        size=34,
        bold=True,
        color=_WHITE,
        center=True,
    )
    _text(
        s,
        Inches(1),
        Inches(2.0),
        W - Inches(2),
        Inches(0.45),
        "AI-Powered RAG Project Health Analysis",
        size=16,
        color=_GREY,
        center=True,
    )
    _text(
        s,
        Inches(1),
        int(H * 0.56) + Inches(0.05),
        W - Inches(2),
        Inches(0.45),
        f"Report Date: {summary.get('report_date', date.today().isoformat())}  |  "
        f"{summary.get('total_projects', 0)} Projects Analysed",
        size=13,
        bold=True,
        color=_WHITE,
        center=True,
    )
    _text(
        s,
        Inches(1),
        H - Inches(0.45),
        W - Inches(2),
        Inches(0.38),
        "AI Project Health Reporting System  |  Confidential",
        size=9,
        color=_WHITE,
        center=True,
    )


def _p_overview(prs, summary):
    from pptx.util import Inches

    s = _new_blank(prs)
    W, H = prs.slide_width, prs.slide_height
    _set_bg(s, *_LIGHT)
    _header_bar(s, prs, "📊  Portfolio Overview")

    dist = summary.get("rag_distribution", {})
    total = summary.get("total_projects", 0)
    avg = summary.get("avg_composite_score", "N/A")
    signal = summary.get("health_signal", "")
    sig_c = (
        _RED
        if "high risk" in signal.lower()
        else (_AMBER if "pressure" in signal.lower() else _GREEN)
    )

    kpis = [
        ("Total", str(total), _DARK),
        ("🟢 Green", str(dist.get("Green", 0)), _GREEN),
        ("🟡 Amber", str(dist.get("Amber", 0)), _AMBER),
        ("🔴 Red", str(dist.get("Red", 0)), _RED),
        ("Avg Score", str(avg), _BLUE),
    ]
    bw, bh, gap, sx, ty = (
        Inches(2.2),
        Inches(1.1),
        Inches(0.22),
        Inches(0.5),
        Inches(0.85),
    )
    for i, (lbl, val, col) in enumerate(kpis):
        x = sx + i * (bw + gap)
        _box(s, x, ty, bw, bh, fill=col)
        _text(
            s,
            x,
            ty + Inches(0.08),
            bw,
            Inches(0.55),
            val,
            size=26,
            bold=True,
            color=_WHITE,
            center=True,
        )
        _text(
            s,
            x,
            ty + Inches(0.65),
            bw,
            Inches(0.35),
            lbl,
            size=10,
            color=_WHITE,
            center=True,
        )

    _box(s, Inches(0.5), Inches(2.15), W - Inches(1), Inches(0.5), fill=sig_c)
    _text(
        s,
        Inches(0.6),
        Inches(2.15),
        W - Inches(1.2),
        Inches(0.5),
        f"  Portfolio Signal: {signal}",
        size=13,
        bold=True,
        color=_WHITE,
    )

    max_c = max(dist.get("Green", 0), dist.get("Amber", 0), dist.get("Red", 0), 1)
    bar_h = Inches(2.4)
    bar_w = Inches(1.8)
    bar_ty = Inches(2.85)
    for i, (rag, col) in enumerate(
        [("Green", _GREEN), ("Amber", _AMBER), ("Red", _RED)]
    ):
        cnt = dist.get(rag, 0)
        bh2 = int(bar_h * cnt / max_c) if max_c else 0
        bx = Inches(2.5 + i * (bar_w + Inches(0.6)))
        if bh2 > 0:
            _box(s, bx, bar_ty + bar_h - bh2, bar_w, bh2, fill=col)
        _text(
            s,
            bx,
            bar_ty + bar_h + Inches(0.06),
            bar_w,
            Inches(0.4),
            f"{rag}  {cnt}",
            size=12,
            bold=True,
            color=_DARK,
            center=True,
        )

    _footer_bar(s, prs, "Portfolio Overview", "Slide 2 / 7")


def _p_scorecard(prs, projects):
    from pptx.util import Inches

    s = _new_blank(prs)
    W, H = prs.slide_width, prs.slide_height
    _set_bg(s, *_LIGHT)
    _header_bar(s, prs, "🏗️  Project Health Scorecard")

    cols = ["Project", "RAG", "Score", "Compl.", "Delayed", "Risks", "Summary"]
    cw = [
        Inches(2.6),
        Inches(1.1),
        Inches(0.9),
        Inches(0.9),
        Inches(0.9),
        Inches(0.8),
        Inches(5.5),
    ]
    row_h = Inches(0.5)
    hdr_y = Inches(0.85)
    x = Inches(0.25)
    for ci, (col, cwidth) in enumerate(zip(cols, cw)):
        _box(s, x, hdr_y, cwidth, row_h, fill=_DARK)
        _text(
            s,
            x,
            hdr_y,
            cwidth,
            row_h,
            col,
            size=10,
            bold=True,
            color=_WHITE,
            center=True,
        )
        x += cwidth

    for ri, p in enumerate(projects[:10]):
        ry = hdr_y + row_h + ri * row_h
        bg = (235, 245, 255) if ri % 2 == 0 else _WHITE
        rag = p.rag_status or "Unknown"
        vals = [
            p.name[:32],
            "",
            str(p.metrics.get("composite_score", "N/A")),
            f"{p.metrics.get('completion_percent','N/A')}%",
            str(p.metrics.get("delayed_tasks", 0)),
            str(p.metrics.get("open_risks", 0)),
            (p.executive_summary or "")[:70],
        ]
        x = Inches(0.25)
        for ci, (val, cwidth) in enumerate(zip(vals, cw)):
            if ci == 1:
                _box(s, x, ry, cwidth, row_h, fill=bg)
                bw2 = Inches(0.8)
                _box(
                    s,
                    x + (cwidth - bw2) // 2,
                    ry + Inches(0.08),
                    bw2,
                    Inches(0.32),
                    fill=_RAG_RGB.get(rag, _GREY),
                )
                _text(
                    s,
                    x + (cwidth - bw2) // 2,
                    ry + Inches(0.08),
                    bw2,
                    Inches(0.32),
                    rag,
                    size=8,
                    bold=True,
                    color=_WHITE,
                    center=True,
                )
            else:
                _box(s, x, ry, cwidth, row_h, fill=bg)
                _text(
                    s,
                    x + Inches(0.05),
                    ry,
                    cwidth - Inches(0.05),
                    row_h,
                    val,
                    size=9,
                    color=_DARK,
                    center=(ci > 1 and ci < 6),
                )
            x += cwidth

    _footer_bar(s, prs, "Project Health Scorecard", "Slide 3 / 7")


def _p_risks(prs, projects):
    from pptx.util import Inches

    s = _new_blank(prs)
    W, H = prs.slide_width, prs.slide_height
    _set_bg(s, *_LIGHT)
    _header_bar(s, prs, "🚨  Key Risks & Issues", bg=(192, 0, 0))

    all_risks = []
    for p in projects:
        for r in (p.identified_risks or [])[:2]:
            all_risks.append((p.name, r, p.rag_status or "Unknown"))
        for r in p.risks:
            if (r.severity or "").lower() in {"high", "critical"}:
                all_risks.append((p.name, r.description or "", "Red"))

    seen, unique = set(), []
    for proj, txt, st in all_risks:
        key = f"{proj}|{txt[:40]}"
        if key not in seen and txt:
            seen.add(key)
            unique.append((proj, txt, st))
        if len(unique) >= 10:
            break

    if not unique:
        _text(
            s,
            Inches(0.5),
            Inches(1.2),
            W - Inches(1),
            Inches(1),
            "No significant risks identified.",
            size=14,
            color=_DARK,
        )
    else:
        row_h = Inches(0.52)
        for i, (proj, txt, status) in enumerate(unique):
            y = Inches(0.88) + i * row_h
            rgb = _RAG_RGB.get(status, _GREY)
            bg = (
                (253, 237, 236)
                if status == "Red"
                else (254, 249, 231) if status == "Amber" else (213, 245, 227)
            )
            _box(s, Inches(0.25), y, W - Inches(0.5), row_h - Inches(0.04), fill=bg)
            _box(s, Inches(0.25), y, Inches(0.1), row_h - Inches(0.04), fill=rgb)
            _text(
                s,
                Inches(0.45),
                y,
                Inches(2.2),
                row_h - Inches(0.04),
                f"[{proj[:22]}]",
                size=9,
                bold=True,
                color=_DARK,
            )
            _text(
                s,
                Inches(2.65),
                y,
                W - Inches(3),
                row_h - Inches(0.04),
                txt[:130],
                size=9,
                color=_DARK,
            )

    _footer_bar(s, prs, "Risks & Issues", "Slide 4 / 7")


def _p_recommendations(prs, projects):
    from pptx.util import Inches

    s = _new_blank(prs)
    W, H = prs.slide_width, prs.slide_height
    _set_bg(s, *_LIGHT)
    _header_bar(s, prs, "✅  Recommendations & Next Steps", bg=_BLUE)

    entries = []
    for p in projects:
        for rec in (p.recommendations or [])[:2] + (p.next_week_priorities or [])[:1]:
            if rec:
                entries.append((p.name, rec))
        if len(entries) >= 12:
            break

    if not entries:
        _text(
            s,
            Inches(0.5),
            Inches(1.2),
            W - Inches(1),
            Inches(1),
            "Enable Gemini AI to generate recommendations.",
            size=14,
            color=_DARK,
        )
    else:
        left_col = [e for i, e in enumerate(entries[:12]) if i % 2 == 0]
        right_col = [e for i, e in enumerate(entries[:12]) if i % 2 == 1]
        row_h = Inches(0.52)
        for ci, col in enumerate([left_col, right_col]):
            x = Inches(0.3) if ci == 0 else Inches(6.9)
            for ri, (proj, rec) in enumerate(col):
                y = Inches(0.88) + ri * row_h
                _box(s, x, y, Inches(6.3), row_h - Inches(0.06), fill=(235, 245, 255))
                _box(s, x, y, Inches(0.1), row_h - Inches(0.06), fill=_BLUE)
                _text(
                    s,
                    x + Inches(0.18),
                    y,
                    Inches(6),
                    Inches(0.22),
                    f"[{proj[:24]}]",
                    size=8,
                    bold=True,
                    color=_BLUE,
                )
                _text(
                    s,
                    x + Inches(0.18),
                    y + Inches(0.22),
                    Inches(6),
                    Inches(0.26),
                    rec[:110],
                    size=9,
                    color=_DARK,
                )

    _footer_bar(s, prs, "Recommendations", "Slide 5 / 7")


def _p_methodology(prs):
    from pptx.util import Inches

    s = _new_blank(prs)
    W, H = prs.slide_width, prs.slide_height
    _set_bg(s, *_LIGHT)
    _header_bar(s, prs, "📘  RAG Scoring Methodology")

    dims = [
        ("Schedule Performance", "25%", "Variance to baseline; delayed task ratio"),
        ("Milestone Completion", "20%", "Overdue milestones vs total"),
        ("Task Completion", "20%", "Completed tasks as % of total"),
        ("Risk Exposure", "15%", "Open risks weighted by severity"),
        ("Critical Path Blockers", "10%", "Critical tasks with positive variance"),
        ("Data Quality", "10%", "Completeness of required plan fields"),
    ]
    cw = [Inches(3.2), Inches(1.0), Inches(8.8)]
    row_h = Inches(0.5)
    hdr_y = Inches(0.88)
    headers = ["Dimension", "Weight", "What It Measures"]
    x = Inches(0.3)
    for hdr, cwidth in zip(headers, cw):
        _box(s, x, hdr_y, cwidth, row_h, fill=_DARK)
        _text(
            s,
            x,
            hdr_y,
            cwidth,
            row_h,
            hdr,
            size=10,
            bold=True,
            color=_WHITE,
            center=True,
        )
        x += cwidth

    for ri, (dim, wt, criteria) in enumerate(dims):
        ry = hdr_y + row_h + ri * row_h
        bg = (235, 245, 255) if ri % 2 == 0 else _WHITE
        x = Inches(0.3)
        for val, cwidth in zip([dim, wt, criteria], cw):
            _box(s, x, ry, cwidth, row_h - Inches(0.04), fill=bg)
            _text(
                s,
                x + Inches(0.08),
                ry,
                cwidth - Inches(0.1),
                row_h - Inches(0.04),
                val,
                size=10,
                color=_DARK,
                center=(cwidth < Inches(1.5)),
            )
            x += cwidth

    ty = hdr_y + row_h * 7 + Inches(0.2)
    for i, (lbl, desc, col) in enumerate(
        [
            ("🟢 GREEN", "Score ≥ 0.70", _GREEN),
            ("🟡 AMBER", "0.40 – 0.69", _AMBER),
            ("🔴 RED", "< 0.40", _RED),
        ]
    ):
        tx = Inches(1.5 + i * Inches(3.6))
        _box(s, tx, ty, Inches(3.2), Inches(0.55), fill=col)
        _text(
            s,
            tx,
            ty,
            Inches(3.2),
            Inches(0.28),
            lbl,
            size=12,
            bold=True,
            color=_WHITE,
            center=True,
        )
        _text(
            s,
            tx,
            ty + Inches(0.28),
            Inches(3.2),
            Inches(0.25),
            desc,
            size=9,
            color=_WHITE,
            center=True,
        )

    _footer_bar(s, prs, "RAG Methodology", "Slide 6 / 7")


def _p_closing(prs, summary, projects):
    from pptx.util import Inches

    s = _new_blank(prs)
    W, H = prs.slide_width, prs.slide_height
    _set_bg(s, *_DARK)
    _box(s, 0, Inches(0.75), W, Inches(0.07), fill=_BLUE)
    _text(
        s,
        Inches(0.5),
        Inches(0.1),
        W - Inches(1),
        Inches(0.65),
        "Executive Summary & Closing",
        size=24,
        bold=True,
        color=_WHITE,
    )
    signal = summary.get("health_signal", "")
    sig_c = (
        _RED
        if "high risk" in signal.lower()
        else (_AMBER if "pressure" in signal.lower() else _GREEN)
    )
    _box(s, Inches(0.5), Inches(1.0), W - Inches(1), Inches(0.55), fill=sig_c)
    _text(
        s,
        Inches(0.6),
        Inches(1.0),
        W - Inches(1.2),
        Inches(0.55),
        f"  {signal}",
        size=13,
        bold=True,
        color=_WHITE,
    )
    y = Inches(1.75)
    for p in projects[:7]:
        rag = p.rag_status or "Unknown"
        sym = {"Green": "🟢", "Amber": "🟡", "Red": "🔴"}.get(rag, "⚪")
        _text(
            s,
            Inches(0.5),
            y,
            W - Inches(1),
            Inches(0.48),
            f"{sym}  {p.name}:  {(p.executive_summary or 'No summary')[:130]}",
            size=10,
            color=_WHITE,
        )
        y += Inches(0.52)
    _box(s, 0, H - Inches(0.5), W, Inches(0.5), fill=_BLUE)
    _text(
        s,
        Inches(0.5),
        H - Inches(0.48),
        W - Inches(1),
        Inches(0.4),
        "AI Project Health Reporting System  •  Powered by Gemini AI  •  Confidential",
        size=9,
        color=_WHITE,
        center=True,
    )


# ══════════════════════════════════════════════════════════════════════════════
# PER-PROJECT SLIDES (5 slides)
# ══════════════════════════════════════════════════════════════════════════════


def _proj_title(prs, project):
    from pptx.util import Inches

    s = _new_blank(prs)
    W, H = prs.slide_width, prs.slide_height
    rag = project.rag_status or "Unknown"
    col = _RAG_RGB.get(rag, _GREY)
    score = project.metrics.get("composite_score", "N/A")
    _set_bg(s, *_DARK)
    _box(s, 0, H - Inches(0.55), W, Inches(0.55), fill=_NAVY)
    _box(s, 0, H - Inches(0.55) - Inches(0.06), W, Inches(0.06), fill=col)
    _text(
        s,
        Inches(1),
        Inches(0.9),
        W - Inches(2),
        Inches(0.45),
        "Weekly Project Health Report",
        size=14,
        color=_GREY,
        center=True,
    )
    _text(
        s,
        Inches(0.5),
        Inches(1.5),
        W - Inches(1),
        Inches(1.0),
        project.name,
        size=36,
        bold=True,
        color=_WHITE,
        center=True,
    )
    _box(s, W // 2 - Inches(1.5), Inches(2.8), Inches(3), Inches(0.55), fill=col)
    _text(
        s,
        W // 2 - Inches(1.5),
        Inches(2.8),
        Inches(3),
        Inches(0.55),
        rag,
        size=16,
        bold=True,
        color=_WHITE,
        center=True,
    )
    _text(
        s,
        Inches(1),
        Inches(3.55),
        W - Inches(2),
        Inches(0.4),
        f"Score: {score}  |  PM: {project.pm or 'N/A'}  |  {date.today().isoformat()}",
        size=12,
        color=_GREY,
        center=True,
    )
    _text(
        s,
        Inches(0.5),
        H - Inches(0.47),
        W - Inches(1),
        Inches(0.4),
        "AI Project Health Reporting System  —  Weekly Report",
        size=9,
        color=_WHITE,
        center=True,
    )


def _proj_dashboard(prs, project):
    from pptx.util import Inches

    s = _new_blank(prs)
    W, H = prs.slide_width, prs.slide_height
    rag = project.rag_status or "Unknown"
    col = _RAG_RGB.get(rag, _GREY)
    m = project.metrics
    conf = f"{project.rag_confidence*100:.0f}%" if project.rag_confidence else "N/A"
    _set_bg(s, *_LIGHT)
    _header_bar(s, prs, f"📈  Health Dashboard — {project.name[:45]}", bg=col)

    kpis = [
        ("Completion", f"{m.get('completion_percent','N/A')}%", _BLUE),
        (
            "Tasks Done",
            f"{m.get('completed_tasks',0)}/{m.get('total_tasks',0)}",
            _GREEN,
        ),
        ("Delayed", str(m.get("delayed_tasks", 0)), _RED),
        ("Overdue MS", str(m.get("overdue_milestones", 0)), _AMBER),
        ("Open Risks", str(m.get("open_risks", 0)), (155, 89, 182)),
        ("Confidence", conf, (26, 188, 156)),
    ]
    bw, bh, gap, sx, ty = (
        Inches(1.9),
        Inches(0.95),
        Inches(0.18),
        Inches(0.35),
        Inches(0.85),
    )
    for i, (lbl, val, c) in enumerate(kpis):
        x = sx + i * (bw + gap)
        _box(s, x, ty, bw, bh, fill=c)
        _text(
            s,
            x,
            ty + Inches(0.1),
            bw,
            Inches(0.55),
            val,
            size=22,
            bold=True,
            color=_WHITE,
            center=True,
        )
        _text(
            s,
            x,
            ty + Inches(0.65),
            bw,
            Inches(0.25),
            lbl,
            size=9,
            color=_WHITE,
            center=True,
        )

    if project.executive_summary:
        _box(
            s,
            Inches(0.3),
            Inches(2.0),
            W - Inches(0.6),
            Inches(0.85),
            fill=(234, 244, 251),
        )
        _text(
            s,
            Inches(0.45),
            Inches(2.0),
            W - Inches(0.9),
            Inches(0.85),
            (project.executive_summary or "")[:220],
            size=10,
            color=_DARK,
        )

    # Dimension bars
    dims = list(m.get("dimensions", {}).items())[:6]
    dy = Inches(3.05)
    bar_max_w = W - Inches(5.5)
    for dim_name, data in dims:
        score = float(data.get("score", 0))
        d_col = _GREEN if score >= 0.7 else (_AMBER if score >= 0.4 else _RED)
        filled = int(bar_max_w * score)
        wt = m.get("weights", {}).get(dim_name, 0)
        rat = (data.get("rationale", "") or "")[:50]
        _text(
            s,
            Inches(0.3),
            dy,
            Inches(2.5),
            Inches(0.35),
            dim_name.replace("_", " ").title(),
            size=9,
            bold=True,
            color=_DARK,
        )
        _box(
            s,
            Inches(2.9),
            dy + Inches(0.05),
            bar_max_w,
            Inches(0.26),
            fill=(220, 220, 220),
        )
        if filled > 0:
            _box(s, Inches(2.9), dy + Inches(0.05), filled, Inches(0.26), fill=d_col)
        _text(
            s,
            Inches(2.95),
            dy + Inches(0.05),
            Inches(0.8),
            Inches(0.26),
            f"{score:.2f}",
            size=8,
            bold=True,
            color=_WHITE,
        )
        _text(
            s,
            Inches(2.9) + bar_max_w + Inches(0.1),
            dy,
            Inches(1.0),
            Inches(0.35),
            f"{wt:.0%}",
            size=8,
            color=_GREY,
        )
        _text(
            s,
            Inches(2.9) + bar_max_w + Inches(1.2),
            dy,
            Inches(2.8),
            Inches(0.35),
            rat,
            size=8,
            color=_GREY,
        )
        dy += Inches(0.42)

    _footer_bar(s, prs, "Health Dashboard", "Slide 2 / 5")


def _proj_progress(prs, project):
    from pptx.util import Inches

    s = _new_blank(prs)
    W, H = prs.slide_width, prs.slide_height
    _set_bg(s, *_LIGHT)
    _header_bar(s, prs, "📋  Progress & Milestone Tracker")

    overdue = project.overdue_milestones
    oy = Inches(0.88)
    if overdue:
        _box(s, Inches(0.25), oy, W - Inches(0.5), Inches(0.35), fill=(231, 76, 60))
        _text(
            s,
            Inches(0.35),
            oy,
            W - Inches(0.7),
            Inches(0.35),
            f"  ⚠  {len(overdue)} Overdue Milestone(s)",
            size=11,
            bold=True,
            color=_WHITE,
        )
        oy += Inches(0.38)
        for ms in overdue[:3]:
            _box(
                s, Inches(0.25), oy, W - Inches(0.5), Inches(0.38), fill=(253, 237, 236)
            )
            _text(
                s,
                Inches(0.35),
                oy,
                W - Inches(0.7),
                Inches(0.38),
                f"  {ms.name[:50]}  |  Due: {ms.planned_finish.isoformat() if ms.planned_finish else 'N/A'}  |  Overdue: +{ms.variance_days:.0f}d",
                size=9,
                color=_DARK,
            )
            oy += Inches(0.4)

    all_tasks = project.tasks + project.milestones
    avail_h = H - oy - Inches(0.6)
    max_rows = max(1, int(avail_h / Inches(0.42)))

    cols = ["Task / Milestone", "Status", "% Done", "MS", "Crit", "Delayed", "Variance"]
    cw = [
        Inches(4.2),
        Inches(1.3),
        Inches(0.9),
        Inches(0.7),
        Inches(0.7),
        Inches(0.85),
        Inches(1.0),
    ]
    row_h = Inches(0.42)
    x = Inches(0.25)
    for col, cwidth in zip(cols, cw):
        _box(s, x, oy, cwidth, row_h, fill=_DARK)
        _text(
            s, x, oy, cwidth, row_h, col, size=9, bold=True, color=_WHITE, center=True
        )
        x += cwidth

    for ri, t in enumerate(all_tasks[:max_rows]):
        ry = oy + row_h + ri * row_h
        bg = (235, 245, 255) if ri % 2 == 0 else _WHITE
        ms = "★" if t.is_milestone else ""
        crt = "⚠" if t.is_critical else ""
        dlyd = "🔴" if t.is_delayed else "🟢"
        var = f"{t.variance_days:+.0f}d" if t.variance_days else ""
        pct = f"{t.percent_complete:.0f}%" if t.percent_complete is not None else ""
        vals = [(t.name or "")[:48], t.status or "", pct, ms, crt, dlyd, var]
        x = Inches(0.25)
        for val, cwidth in zip(vals, cw):
            _box(s, x, ry, cwidth, row_h, fill=bg)
            _text(
                s,
                x + Inches(0.06),
                ry,
                cwidth - Inches(0.06),
                row_h,
                val,
                size=8,
                color=_DARK,
                center=(cwidth < Inches(1.2)),
            )
            x += cwidth

    _footer_bar(s, prs, "Progress & Milestones", "Slide 3 / 5")


def _proj_risks_slide(prs, project):
    from pptx.util import Inches

    s = _new_blank(prs)
    W, H = prs.slide_width, prs.slide_height
    _set_bg(s, *_LIGHT)
    _header_bar(s, prs, "🚨  Risk Register", bg=(192, 0, 0))

    all_risks = list(project.risks[:12])
    id_risks = project.identified_risks or []
    row_h = Inches(0.46)
    oy = Inches(0.88)

    if all_risks:
        cols = ["#", "Description", "Severity", "Probability", "Owner", "Status"]
        cw = [
            Inches(0.5),
            Inches(5.8),
            Inches(1.1),
            Inches(1.2),
            Inches(2.2),
            Inches(1.5),
        ]
        x = Inches(0.25)
        for col, cwidth in zip(cols, cw):
            _box(s, x, oy, cwidth, row_h, fill=(192, 0, 0))
            _text(
                s,
                x,
                oy,
                cwidth,
                row_h,
                col,
                size=9,
                bold=True,
                color=_WHITE,
                center=True,
            )
            x += cwidth
        for ri, r in enumerate(all_risks):
            ry = oy + row_h + ri * row_h
            bg = (
                (253, 237, 236)
                if (r.severity or "").lower() in {"high", "critical"}
                else (
                    (254, 249, 231)
                    if (r.severity or "").lower() == "medium"
                    else (213, 245, 227)
                )
            )
            vals = [
                str(ri + 1),
                (r.description or "")[:70],
                r.severity or "",
                r.probability or "",
                (r.owner or "")[:25],
                r.status or "",
            ]
            x = Inches(0.25)
            for val, cwidth in zip(vals, cw):
                _box(s, x, ry, cwidth, row_h, fill=bg)
                _text(
                    s,
                    x + Inches(0.05),
                    ry,
                    cwidth - Inches(0.05),
                    row_h,
                    val,
                    size=8,
                    color=_DARK,
                    center=(cwidth < Inches(1.5)),
                )
                x += cwidth
    elif id_risks:
        for i, r in enumerate(id_risks[:10]):
            ry = oy + i * Inches(0.5)
            _box(
                s, Inches(0.25), ry, W - Inches(0.5), Inches(0.44), fill=(253, 237, 236)
            )
            _box(s, Inches(0.25), ry, Inches(0.1), Inches(0.44), fill=_RED)
            _text(
                s,
                Inches(0.42),
                ry,
                W - Inches(0.7),
                Inches(0.44),
                r[:130],
                size=10,
                color=_DARK,
            )
    else:
        _text(
            s,
            Inches(0.5),
            Inches(1.5),
            W - Inches(1),
            Inches(1),
            "No risks identified.",
            size=14,
            color=_DARK,
        )

    _footer_bar(s, prs, "Risk Register", "Slide 4 / 5")


def _proj_recs_slide(prs, project):
    from pptx.util import Inches

    s = _new_blank(prs)
    W, H = prs.slide_width, prs.slide_height
    _set_bg(s, *_LIGHT)
    _header_bar(s, prs, "✅  Recommendations & Next Steps", bg=_BLUE)

    recs = project.recommendations or []
    prios = project.next_week_priorities or []
    missing = project.missing_data or []
    m = project.metrics

    col_w = (W - Inches(0.9)) // 2
    _text(
        s,
        Inches(0.3),
        Inches(0.85),
        col_w,
        Inches(0.35),
        "✅  Recommendations",
        size=12,
        bold=True,
        color=_DARK,
    )
    _text(
        s,
        Inches(0.45) + col_w,
        Inches(0.85),
        col_w,
        Inches(0.35),
        "📅  Next Week Priorities",
        size=12,
        bold=True,
        color=_DARK,
    )

    for i, rec in enumerate(recs[:8]):
        y = Inches(1.25) + i * Inches(0.52)
        _box(s, Inches(0.3), y, col_w, Inches(0.48), fill=(234, 244, 251))
        _box(s, Inches(0.3), y, Inches(0.1), Inches(0.48), fill=_BLUE)
        _text(
            s,
            Inches(0.48),
            y,
            col_w - Inches(0.25),
            Inches(0.48),
            rec[:120],
            size=9,
            color=_DARK,
        )

    for i, p in enumerate(prios[:8]):
        y = Inches(1.25) + i * Inches(0.52)
        _box(s, Inches(0.45) + col_w, y, col_w, Inches(0.48), fill=(254, 249, 231))
        _box(s, Inches(0.45) + col_w, y, Inches(0.1), Inches(0.48), fill=_AMBER)
        _text(
            s,
            Inches(0.63) + col_w,
            y,
            col_w - Inches(0.25),
            Inches(0.48),
            p[:120],
            size=9,
            color=_DARK,
        )

    if missing:
        by = H - Inches(1.6)
        _box(s, Inches(0.3), by, W - Inches(0.6), Inches(0.32), fill=(231, 76, 60))
        _text(
            s,
            Inches(0.4),
            by,
            W - Inches(0.8),
            Inches(0.32),
            "  ⚠  Missing / Low Quality Data",
            size=10,
            bold=True,
            color=_WHITE,
        )
        for i, d in enumerate(missing[:2]):
            _text(
                s,
                Inches(0.4),
                by + Inches(0.34) + i * Inches(0.3),
                W - Inches(0.8),
                Inches(0.28),
                f"• {d[:130]}",
                size=9,
                color=_DARK,
            )

    _footer_bar(s, prs, "Recommendations & Next Steps", "Slide 5 / 5")


# ══════════════════════════════════════════════════════════════════════════════
# Public API
# ══════════════════════════════════════════════════════════════════════════════


def _save_prs(prs, path: Path) -> None:
    buf = io.BytesIO()
    prs.save(buf)
    path.write_bytes(buf.getvalue())


def generate_pptx(
    projects: list[Project],
    portfolio_summary: dict[str, Any],
    output_dir: str | Path,
) -> Path:
    """Build a 7-slide portfolio PPTX."""
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)
    path = out / "monthly_portfolio_presentation.pptx"
    prs = _prs()
    _p_title(prs, portfolio_summary)
    _p_overview(prs, portfolio_summary)
    _p_scorecard(prs, projects)
    _p_risks(prs, projects)
    _p_recommendations(prs, projects)
    _p_methodology(prs)
    _p_closing(prs, portfolio_summary, projects)
    _save_prs(prs, path)
    logger.info("Portfolio PPTX saved: %s", path)
    return path


def generate_project_pptx(project: Project, output_dir: str | Path) -> Path:
    """Build a 5-slide per-project PPTX weekly report."""
    import re

    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)
    safe = re.sub(r"[^\w\-]", "_", project.name)
    path = out / f"{safe}_weekly_presentation.pptx"
    prs = _prs()
    _proj_title(prs, project)
    _proj_dashboard(prs, project)
    _proj_progress(prs, project)
    _proj_risks_slide(prs, project)
    _proj_recs_slide(prs, project)
    _save_prs(prs, path)
    logger.info("Project PPTX saved: %s", path)
    return path


def generate_pptx_bytes(projects, portfolio_summary) -> bytes:
    """Return portfolio PPTX as raw bytes."""
    prs = _prs()
    _p_title(prs, portfolio_summary)
    _p_overview(prs, portfolio_summary)
    _p_scorecard(prs, projects)
    _p_risks(prs, projects)
    _p_recommendations(prs, projects)
    _p_methodology(prs)
    _p_closing(prs, portfolio_summary, projects)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
